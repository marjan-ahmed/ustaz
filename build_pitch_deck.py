#!/usr/bin/env python3
"""Generate the USTAZ pitch deck (.pptx) , NIC Karachi panel-ready v2."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ORANGE = RGBColor(0xDB, 0x4B, 0x0D)
NAVY   = RGBColor(0x11, 0x18, 0x28)
CREAM  = RGBColor(0xFC, 0xEE, 0xE2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x3A, 0x3F, 0x4A)
LGRAY  = RGBColor(0x6B, 0x70, 0x7A)
LITEOR = RGBColor(0xF8, 0xC9, 0xA8)
GREEN  = RGBColor(0x1E, 0x9E, 0x5A)
RED    = RGBColor(0xC2, 0x3B, 0x22)
AMB    = RGBColor(0xC9, 0x7A, 0x12)
CL     = RGBColor(0xCF, 0xD3, 0xDA)

FONT = "Calibri"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def oval(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=15, color=GRAY, gap=8, lead_color=ORANGE):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        lead = p.add_run(); lead.text = "\u25b8  "
        lead.font.size = Pt(size); lead.font.bold = True
        lead.font.color.rgb = lead_color; lead.font.name = FONT
        if isinstance(it, tuple):
            head, rest = it
            r1 = p.add_run(); r1.text = head
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = NAVY; r1.font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.bold = False
            r2.font.color.rgb = color; r2.font.name = FONT
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


def footer(s, n):
    text(s, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
         [[("USTAZ", 9, True, ORANGE), ("   ustaz-bice.vercel.app", 9, False, LGRAY)]])
    text(s, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.35),
         [[(str(n), 9, True, LGRAY)]], align=PP_ALIGN.RIGHT)


def content_slide(no, kicker, title, body_fn):
    s = slide()
    rect(s, 0, 0, EMU_W, EMU_H, WHITE)
    rect(s, Inches(0.5), Inches(0.7), Inches(0.12), Inches(1.15), ORANGE)
    oval(s, Inches(12.2), Inches(-0.6), Inches(1.8), Inches(1.8), CREAM)
    text(s, Inches(0.8), Inches(0.62), Inches(10), Inches(0.4),
         [[(kicker, 12, True, ORANGE)]])
    text(s, Inches(0.78), Inches(1.0), Inches(11.5), Inches(0.9),
         [[(title, 30, True, NAVY)]])
    body_fn(s)
    footer(s, no)
    return s


# =====================================================================
# 1 , COVER (with traction bar)
# =====================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, WHITE)
rect(s, Inches(8.9), 0, Inches(4.43), EMU_H, NAVY)
oval(s, Inches(11.0), Inches(2.4), Inches(2.7), Inches(2.7), ORANGE)
oval(s, Inches(10.2), Inches(4.6), Inches(1.2), Inches(1.2), RGBColor(0x1d,0x26,0x3a))
rect(s, Inches(0.9), Inches(2.5), Inches(0.14), Inches(2.0), ORANGE)
text(s, Inches(1.25), Inches(2.45), Inches(7), Inches(2.2),
     [[("USTAZ", 66, True, NAVY)],
      [("Trusted home-service professionals, on demand.", 20, True, ORANGE)]],
     space_after=10)
text(s, Inches(1.28), Inches(4.55), Inches(7), Inches(1.0),
     [[("Pakistan's verified marketplace for electricians, plumbers, carpenters,", 14, False, GRAY)],
      [("AC, solar & cleaning , with live tracking and a 3-day work guarantee.", 14, False, GRAY)]],
     space_after=4)
# Traction bar
card_y = Inches(5.3); card_w = Inches(2.2); card_h = Inches(0.7)
metrics_cover = [("6", "Service categories live"), ("8+", "Core systems shipped"), ("100%", "Tables RLS-secured")]
cx = Inches(1.28)
for val, label in metrics_cover:
    rect(s, cx, card_y, card_w, card_h, CREAM)
    text(s, cx + Inches(0.1), card_y, card_w, card_h,
         [[(val, 16, True, ORANGE), ("  " + label, 9.5, False, GRAY)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    cx += card_w + Inches(0.15)
text(s, Inches(1.28), Inches(6.15), Inches(7), Inches(0.5),
     [[("\u25cf Live MVP:  ", 13, True, ORANGE), ("ustaz-bice.vercel.app", 13, True, NAVY)]])
text(s, Inches(8.9), Inches(5.6), Inches(4.43), Inches(0.5),
     [[("Applying for NIC Karachi incubation", 11, False, CL)]],
     align=PP_ALIGN.CENTER)

# =====================================================================
# 2 , PROBLEM (moved up , Problem before Solution per NIC framework)
# =====================================================================
def s2(s):
    text(s, Inches(0.8), Inches(1.9), Inches(11.6), Inches(0.8),
         [[("The anxiety: ", 14, True, ORANGE),
           ("a Karachi mother's AC dies in the June heat. She must let an unknown man into her home, with no idea who he is, what he'll charge, when he'll arrive, or whether it'll break again tomorrow.", 13, False, GRAY)]],
         line_spacing=1.12)
    rows = [
        ("No verification", "Support", "An unknown stranger in your home , no identity check, no accountability."),
        ("Opaque pricing", "Financial", "Overcharging is routine , no standard rates, no upfront quote, no receipt."),
        ("No recourse / guarantee", "Process", "If the work fails a day later, the worker vanishes , broken workflow."),
        ("No visibility / ETA", "Productivity", '"He\'s on the way" means waiting blind for hours , like a ride with no map.'),
        ("Supply side ignored", "Financial + Process", "Skilled workers have no steady demand, no digital reach, no portable reputation."),
    ]
    y = Inches(2.85)
    for pain, cat, desc in rows:
        rect(s, Inches(0.8), y, Inches(0.12), Inches(0.6), ORANGE)
        text(s, Inches(1.05), y, Inches(3.2), Inches(0.6), [[(pain, 13, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
        rect(s, Inches(4.35), y + Inches(0.08), Inches(1.95), Inches(0.44), CREAM)
        text(s, Inches(4.35), y + Inches(0.08), Inches(1.95), Inches(0.44),
             [[(cat, 10.5, True, ORANGE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(6.5), y, Inches(6.0), Inches(0.6), [[(desc, 11.5, False, GRAY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += Inches(0.66)
    text(s, Inches(0.8), Inches(6.3), Inches(11.6), Inches(0.5),
         [[("Pain mapped to NIC's four categories , Financial \u00b7 Productivity \u00b7 Process \u00b7 Support.", 11, True, NAVY)]])
content_slide(2, "1. PROBLEM & CUSTOMER PAIN POINTS", "Home services in Pakistan are broken on trust", s2)

# =====================================================================
# 3 , UVP (moved after Problem , connects pain to promise)
# =====================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, CREAM)
rect(s, 0, Inches(2.6), EMU_W, Inches(2.3), NAVY)
oval(s, Inches(11.3), Inches(0.5), Inches(1.6), Inches(1.6), ORANGE)
text(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5), [[("2. UNIQUE VALUE PROPOSITION", 12, True, ORANGE)]])
text(s, Inches(0.8), Inches(1.45), Inches(11.5), Inches(0.9),
     [[("Trust a stranger in your home , for the first time.", 28, True, NAVY)]])
text(s, Inches(1.0), Inches(2.85), Inches(11.3), Inches(1.9),
     [[("USTAZ phone-verifies every pro, shows you exactly where they are, and ", 18, False, WHITE),
       ("guarantees the work for 3 days , or we fix it free.", 18, True, ORANGE)],
      [("Halal by design , no interest, no hidden charges, just a transparent service fee.", 16, True, LITEOR)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=10, line_spacing=1.16)
chips = ["Feel safe at the door", "See them on the way", "Covered for 3 days", "Halal, transparent fee"]
x = Inches(0.8)
for c in chips:
    w = Inches(2.95)
    rect(s, x, Inches(5.4), w, Inches(0.8), WHITE, line=LITEOR)
    text(s, x, Inches(5.4), w, Inches(0.8), [[(c, 13, True, NAVY)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += Inches(3.05)
footer(s, 3)

# =====================================================================
# 4 , SOLUTION & MVP (with product screenshot placeholder)
# =====================================================================
def s4(s):
    text(s, Inches(0.8), Inches(1.95), Inches(11.5), Inches(0.5),
         [[("An on-demand marketplace of phone-verified pros, matched by location and tracked live.", 16, True, ORANGE)]])
    steps = [("1", "Pick a service & location", "Choose from electrician, plumber, carpenter, AC, solar or cleaning , set your spot on the map."),
             ("2", "Get matched instantly", "PostGIS finds nearby phone-verified providers; they accept in real time."),
             ("3", "Track live & pay on done", "Watch your pro arrive on a live map, chat in-app, pay on completion , covered by a 3-day warranty.")]
    x = Inches(0.8)
    for num, head, body in steps:
        rect(s, x, Inches(2.7), Inches(3.7), Inches(2.5), CREAM)
        oval(s, x + Inches(0.25), Inches(2.95), Inches(0.7), Inches(0.7), ORANGE)
        text(s, x + Inches(0.25), Inches(3.0), Inches(0.7), Inches(0.6),
             [[(num, 22, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(0.25), Inches(3.85), Inches(3.2), Inches(1.3),
             [[(head, 14.5, True, NAVY)], [(body, 11.5, False, GRAY)]], space_after=6, line_spacing=1.08)
        x += Inches(3.95)
    # Product screenshot placeholder
    rect(s, Inches(0.8), Inches(5.35), Inches(5.5), Inches(1.3), NAVY)
    text(s, Inches(0.8), Inches(5.35), Inches(5.5), Inches(1.3),
         [[("[  PRODUCT SCREENSHOT  ]", 14, True, ORANGE)],
          [("Insert: booking flow / live GPS / provider dashboard", 10, False, CL)],
          [("ustaz-bice.vercel.app , live MVP, verifiable by panel", 10, False, CL)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.6), Inches(5.35), Inches(6.0), Inches(1.3),
         [[("Trust built in:  ", 13, True, NAVY)],
          [("Phone-verified providers \u00b7 live GPS \u00b7 in-app chat & push", 11.5, False, GRAY)],
          [("3-day free re-fix \u00b7 two-way ratings \u00b7 prepaid escrow wallet", 11.5, False, GRAY)],
          [("CNIC & background checks launching in Phase 2", 11, True, ORANGE)]],
         space_after=4, line_spacing=1.08)
content_slide(4, "3. SOLUTION & MVP", "Verified. Tracked. Guaranteed.", s4)

# =====================================================================
# 5 , MARKET SIZING (all PKR, Year 3 SOM added)
# =====================================================================
def s5(s):
    data = [("01", "TAM", "~PKR 800 to 1,300B",
             "Pakistan home services (informal + formal). Source: SMEDA / ILO, 2024."),
            ("02", "SAM", "~PKR 160B",
             "Smartphone-enabled urban households , Karachi, Lahore, Islamabad, Faisalabad. Source: PTA / PBS HIES, 2025."),
            ("03", "SOM Y1", "~PKR 1.4M",
             "Year 1 Karachi pilot: 40 providers x 12 jobs/mo x Rs. 72 commission."),
            ("04", "SOM Y3", "~PKR 25M",
             "Year 3: 175 providers x 25 jobs/mo x Rs. 72 x 12 mo , Karachi mature + Lahore entry.")]
    x = Inches(0.8)
    for num, code, figure, desc in data:
        rect(s, x, Inches(2.1), Inches(2.85), Inches(3.7), WHITE, line=LITEOR)
        oval(s, x + Inches(0.95), Inches(1.75), Inches(0.95), Inches(0.95), ORANGE)
        text(s, x + Inches(0.95), Inches(1.82), Inches(0.95), Inches(0.85),
             [[(num, 18, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x, Inches(2.8), Inches(2.85), Inches(0.45),
             [[(code, 14, True, NAVY)]], align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.18), Inches(2.85), Inches(0.6),
             [[(figure, 22, True, ORANGE)]], align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.15), Inches(3.85), Inches(2.55), Inches(1.8),
             [[(desc, 9.5, False, GRAY)]], align=PP_ALIGN.CENTER, line_spacing=1.08)
        x += Inches(3.05)
    rect(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.95), NAVY)
    text(s, Inches(1.05), Inches(6.03), Inches(11.3), Inches(0.85),
         [[("Bottom-up SOM:  ", 12, True, ORANGE),
           ("Y1: 40 providers x 12 jobs/mo x Rs. 600 avg x 12% = ~Rs. 1.4M.   Y3: 175 x 25 jobs/mo x Rs. 600 x 12% = ~Rs. 3.8M/quarter.", 11, False, WHITE)],
          [("Sources: SMEDA/ILO 2024. Y1 is conservative pilot-year only , multi-year ramp to Rs. 25M+.", 10, True, CL)]],
         space_after=3, line_spacing=1.04)
content_slide(5, "4. MARKET SIZING  (TAM \u00b7 SAM \u00b7 SOM)", "A vast, underserved, digitizing market", s5)

# =====================================================================
# 6 , CUSTOMER SEGMENTATION (with acquisition channels)
# =====================================================================
def s6(s):
    rect(s, Inches(0.8), Inches(1.95), Inches(5.65), Inches(4.3), CREAM)
    rect(s, Inches(0.8), Inches(1.95), Inches(5.65), Inches(0.12), ORANGE)
    text(s, Inches(1.1), Inches(2.2), Inches(5.1), Inches(0.6),
         [[('DEMAND , "Safe Household Manager"', 13.5, True, ORANGE)]])
    text(s, Inches(1.1), Inches(2.8), Inches(5.15), Inches(3.3),
         [[("Demographic: ", 12, True, NAVY), ("age 28 to 50, urban Karachi, SEC A/B household.", 12, False, GRAY)],
          [("Geographic: ", 12, True, NAVY), ("DHA, Clifton, Gulshan, PECHS, North Nazimabad.", 12, False, GRAY)],
          [("Psychographic: ", 12, True, NAVY), ("values safety over price; has been overcharged; manages the home; won't risk an unverified stranger; wants a visible ETA and accountability.", 12, False, GRAY)],
          [("Trigger: ", 12, True, ORANGE), ("AC breaks in June , needs a trusted technician today, not a stranger from a Facebook group.", 12, False, GRAY)],
          [("Channel: ", 12, True, NAVY), ("Meta/TikTok ads (Karachi women 25 to 45) + residential society B2B2C partnerships.", 12, False, GRAY)]],
         space_after=7, line_spacing=1.08)
    rect(s, Inches(6.85), Inches(1.95), Inches(5.65), Inches(4.3), NAVY)
    rect(s, Inches(6.85), Inches(1.95), Inches(5.65), Inches(0.12), ORANGE)
    text(s, Inches(7.15), Inches(2.2), Inches(5.1), Inches(0.6),
         [[('SUPPLY , "Underserved Tradesman"', 13.5, True, ORANGE)]])
    text(s, Inches(7.15), Inches(2.8), Inches(5.15), Inches(3.3),
         [[("Demographic: ", 12, True, WHITE), ("age 22 to 45, skilled electrician / plumber / AC tech.", 12, False, CL)],
          [("Geographic: ", 12, True, WHITE), ("works across Karachi, based in lower-income areas.", 12, False, CL)],
          [("Psychographic: ", 12, True, WHITE), ("wants steady, predictable digital demand and a reputation that travels with him; tired of middlemen taking large cuts.", 12, False, CL)],
          [("Trigger: ", 12, True, ORANGE), ("wants to stop depending on ustads / contractors and own his own customer pipeline.", 12, False, CL)],
          [("Channel: ", 12, True, WHITE), ("door-to-door at trade hubs (Saddar, Lea Market) + guild batch-onboarding.", 12, False, CL)]],
         space_after=7, line_spacing=1.08)
    text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         [[("This two-sided pool , Karachi's middle-class households x the city's skilled tradesmen , is our SAM foundation (sizing on previous slide).", 11, False, NAVY)]])
content_slide(6, "5. CUSTOMER SEGMENTATION & PERSONA", "Two personas, served on both sides", s6)

# =====================================================================
# 7 , COMPETITOR ANALYSIS (Urban Company added)
# =====================================================================
def s7(s):
    headers = ["Alternative", "Verified", "Live GPS", "Warranty", "Escrow", "Reputation", "Price clarity"]
    data = [
        ("USTAZ",            "\u2713", "\u2713", "\u2713", "\u2713", "\u2713", "\u2713"),
        ("Urban Company (IN)","\u2713", "\u2713", "\u2713", "\u2713", "\u2713", "\u2713"),
        ("Mahir",            "~", "\u2717", "\u2717", "\u2717", "~", "~"),
        ("RepairGuru",       "~", "\u2717", "\u2717", "\u2717", "~", "~"),
        ("OLX / classifieds","\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717"),
        ("Facebook groups",  "\u2717", "\u2717", "\u2717", "\u2717", "\u2717", "\u2717"),
        ("Word-of-mouth",    "\u2717", "\u2717", "\u2717", "\u2717", "~", "\u2717"),
    ]
    rows = len(data) + 1
    cols = len(headers)
    tbl = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.95), Inches(11.7), Inches(3.55)).table
    tbl.columns[0].width = Inches(2.7)
    for c in range(1, cols):
        tbl.columns[c].width = Inches(1.5)
    for c in range(cols):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
        r = p.add_run(); r.text = headers[c]; r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for ri, row in enumerate(data, start=1):
        is_ustaz = ri == 1
        is_uc = ri == 2
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.fill.solid()
            cell.fill.fore_color.rgb = (LITEOR if is_ustaz else (RGBColor(0xE8,0xF0,0xFE) if is_uc else (CREAM if ri % 2 else WHITE))) if ci else (ORANGE if is_ustaz else (GREEN if is_uc else WHITE))
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            if ci == 0:
                r.text = val; r.font.bold = True
                r.font.color.rgb = WHITE if is_ustaz else (GREEN if is_uc else NAVY); r.font.size = Pt(11.5)
            else:
                r.text = val; r.font.bold = True; r.font.size = Pt(13)
                r.font.color.rgb = GREEN if val == "\u2713" else (AMB if val == "~" else RED)
            r.font.name = FONT
    rect(s, Inches(0.8), Inches(5.62), Inches(11.7), Inches(0.46), CREAM)
    text(s, Inches(1.0), Inches(5.62), Inches(11.3), Inches(0.46),
         [[("The real competitor is informality , ~90 to 95% of the market is still cash, unverified, untracked.  ", 11, True, NAVY),
           ("(industry estimate)", 9, False, LGRAY)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.8), Inches(6.16), Inches(11.7), Inches(0.74), NAVY)
    text(s, Inches(1.0), Inches(6.16), Inches(11.3), Inches(0.74),
         [[("Urban Company validated this model at $2.8B in India , Pakistan's market is structurally identical and 3 years behind.  ", 11, True, ORANGE),
           ("Why bookings stay on-platform (anti disintermediation): 3-day warranty & dispute cover apply only to in-app jobs.", 10, False, WHITE)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=1)
content_slide(7, "6. COMPETITOR ANALYSIS / ALTERNATIVES", "We compete with chaos , and beat it", s7)

# =====================================================================
# 8 , MARKETING & GTM (with cold start section)
# =====================================================================
def s8(s):
    rect(s, Inches(0.8), Inches(1.95), Inches(5.65), Inches(3.85), CREAM)
    rect(s, Inches(0.8), Inches(1.95), Inches(5.65), Inches(0.12), ORANGE)
    text(s, Inches(1.1), Inches(2.2), Inches(5.1), Inches(0.5), [[("DEMAND  \u00b7  acquiring customers", 13, True, ORANGE)]])
    bullets(s, Inches(1.1), Inches(2.8), Inches(5.15), Inches(2.9), [
        ("Meta / TikTok ads , ", "Karachi women 25 to 45, home-owner interest."),
        ("SEO landing pages , ", "per service x city (already shipped , a live channel asset)."),
        ("Residential societies (B2B2C) , ", "one manager unlocks 200+ households."),
        ("Referral loop , ", "customers earn for each new customer."),
        ("Community & mosque networks , ", "trust-building in conservative areas."),
    ], size=12, gap=7)
    rect(s, Inches(6.85), Inches(1.95), Inches(5.65), Inches(3.85), NAVY)
    rect(s, Inches(6.85), Inches(1.95), Inches(5.65), Inches(0.12), ORANGE)
    text(s, Inches(7.15), Inches(2.2), Inches(5.1), Inches(0.5), [[("SUPPLY  \u00b7  acquiring providers", 13, True, ORANGE)]])
    bullets(s, Inches(7.15), Inches(2.8), Inches(5.15), Inches(2.9), [
        ("Door-to-door , ", "Karachi trade hubs (Saddar, Lea Market, Sohrab Goth)."),
        ("Ustad / contractor referrals , ", "master tradesmen refer apprentices."),
        ("Provider-led word-of-mouth , ", "satisfied pros recruit peers, at zero cost."),
        ("Trade associations , ", "electrician & plumber guild outreach."),
    ], size=12, gap=8, color=CL)
    # Cold start strategy
    rect(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), NAVY)
    text(s, Inches(1.05), Inches(5.95), Inches(11.3), Inches(0.9),
         [[("Cold start strategy:  ", 12, True, ORANGE),
           ("We seed supply FIRST through guild batch-onboarding (20 to 30 tradesmen at a time via trade associations), then activate demand with a subsidized first job incentive.", 11, False, WHITE)],
          [("Referral loops on both sides drive CAC down over time as the network compounds.", 10, True, CL)]],
         space_after=3, line_spacing=1.04)
content_slide(8, "7. MARKETING & DISTRIBUTION CHANNELS", "Two acquisition engines, one flywheel", s8)

# =====================================================================
# 9 , REVENUE MODEL (with monetization timeline)
# =====================================================================
def s9(s):
    rect(s, Inches(0.8), Inches(1.95), Inches(6.2), Inches(3.6), NAVY)
    rect(s, Inches(0.8), Inches(1.95), Inches(6.2), Inches(0.12), ORANGE)
    text(s, Inches(1.1), Inches(2.2), Inches(5.7), Inches(0.5), [[("HOW WE EARN  \u00b7  prepaid wallet", 13, True, ORANGE)]])
    steps = [
        ("1", "Free Rs. 500 welcome credit", "Every new provider claims Rs. 500 on their dashboard , start earning with no upfront cost."),
        ("2", "12% deducted per completed job", "On each finished job the platform deducts a 12% commission from the provider's wallet (e.g. a Rs. 500 job = Rs. 60)."),
        ("3", "Top up when credit runs low", "The free Rs. 500 credit covers the first several jobs; then the provider tops up to keep receiving jobs (0 balance = paused)."),
    ]
    y = Inches(2.75)
    for n, head, body in steps:
        oval(s, Inches(1.1), y, Inches(0.45), Inches(0.45), ORANGE)
        text(s, Inches(1.1), y, Inches(0.45), Inches(0.45), [[(n, 13, True, WHITE)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.7), y - Inches(0.04), Inches(5.1), Inches(1.0),
             [[(head, 12.5, True, WHITE)], [(body, 10.5, False, CL)]], space_after=2, line_spacing=1.03)
        y += Inches(0.92)
    rect(s, Inches(7.3), Inches(1.95), Inches(5.2), Inches(3.6), WHITE, line=LITEOR)
    rect(s, Inches(7.3), Inches(1.95), Inches(5.2), Inches(0.12), ORANGE)
    text(s, Inches(7.6), Inches(2.2), Inches(4.6), Inches(0.5), [[("TOP UP PACKAGES", 13, True, ORANGE)]])
    pkgs = [("Starter", "Rs. 500", "Just testing"),
            ("Standard", "Rs. 1,000", "Regular work  \u00b7  popular"),
            ("Pro", "Rs. 2,000", "Full time")]
    y = Inches(2.75)
    for label, amt, tag in pkgs:
        rect(s, Inches(7.6), y, Inches(4.6), Inches(0.78), CREAM)
        rect(s, Inches(7.6), y, Inches(0.1), Inches(0.78), ORANGE)
        text(s, Inches(7.8), y, Inches(2.3), Inches(0.78),
             [[(label, 13, True, NAVY)], [(tag, 9.5, False, GRAY)]], space_after=1, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        text(s, Inches(10.0), y, Inches(2.1), Inches(0.78), [[(amt, 16, True, ORANGE)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.88)
    text(s, Inches(7.6), Inches(5.3), Inches(4.7), Inches(0.4),
         [[("Pay via Easypaisa \u00b7 JazzCash \u00b7 Bank \u2192 upload receipt \u2192 admin credits wallet.", 10, False, GRAY)]], line_spacing=1.0)
    # Monetization timeline
    rect(s, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.25), CREAM)
    text(s, Inches(1.05), Inches(5.7), Inches(11.4), Inches(0.4),
         [[("Monetization timeline:", 12, True, ORANGE)]], space_after=2)
    phases = [("Phase 1", "Free credit seeding", "Providers onboard at zero cost, build habits"),
              ("Phase 2", "First top-ups", "Providers convert to paid as free credit exhausts"),
              ("Phase 3", "Commission flowing", "12% per job, recurring revenue compounds"),
              ("Phase 4", "Break even", "~280 jobs/mo covers cloud + acquisition burn")]
    px = Inches(1.05)
    for label, title, desc in phases:
        rect(s, px, Inches(6.1), Inches(0.08), Inches(0.6), ORANGE)
        text(s, px + Inches(0.15), Inches(6.05), Inches(2.6), Inches(0.7),
             [[(label, 10, True, ORANGE)], [(title, 10, True, NAVY)], [(desc, 8.5, False, GRAY)]],
             space_after=1, line_spacing=1.0)
        px += Inches(2.85)
    text(s, Inches(1.05), Inches(6.85), Inches(11.4), Inches(0.3),
         [[("Sharia (ujrah): fee charged ONLY on completed, customer confirmed jobs; unused prepaid credit fully refundable. No interest, no ambiguity.", 9.5, False, GRAY)]])
content_slide(9, "8. REVENUE STREAMS & BUSINESS MODEL", "Free credit, then a 12% prepaid commission", s9)

# =====================================================================
# 10 , COST STRUCTURE (with founder salaries)
# =====================================================================
def s10(s):
    rect(s, Inches(0.8), Inches(1.95), Inches(3.5), Inches(3.5), CREAM)
    rect(s, Inches(0.8), Inches(1.95), Inches(3.5), Inches(0.12), ORANGE)
    text(s, Inches(1.05), Inches(2.18), Inches(3.0), Inches(0.5), [[("CAPEX \u00b7 one time", 13, True, ORANGE)]])
    bullets(s, Inches(1.05), Inches(2.72), Inches(3.1), Inches(2.6), [
        "Product dev , done (founder-built)",
        "Office space & setup",
        "Brand & design system",
        "Native mobile app build",
        "Verification system (one time build)",
    ], size=11.5, gap=10)
    text(s, Inches(4.55), Inches(2.0), Inches(8), Inches(0.4),
         [[("OPEX \u00b7 monthly  ", 12, True, ORANGE), ("(PKR , illustrative 2025 list prices)", 10, False, GRAY)]])
    rows, cols = 8, 3
    tbl = s.shapes.add_table(rows, cols, Inches(4.55), Inches(2.45), Inches(7.95), Inches(3.2)).table
    tbl.columns[0].width = Inches(3.55); tbl.columns[1].width = Inches(2.2); tbl.columns[2].width = Inches(2.2)
    hdr = ["Line item", "Pre-revenue", "500 jobs/mo"]
    data = [
        ("Supabase + Vercel", "~8K to 12K", "~15K to 25K"),
        ("Twilio OTP", "~3K", "~8K"),
        ("Google Maps API", "~2K", "~6K"),
        ("FCM / Resend", "Free", "Free"),
        ("Domain & misc", "~3K", "~3K"),
        ("Founder salaries", "Deferred (equity)", "Deferred (equity)"),
        ("Total cloud burn", "~16K to 20K", "~32K to 42K"),
    ]
    for c in range(cols):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr[c]; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    for ri, row in enumerate(data, start=1):
        total = ri == len(data)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.fill.solid()
            cell.fill.fore_color.rgb = LITEOR if total else (CREAM if ri % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val; r.font.size = Pt(11); r.font.bold = (ci == 0 or total)
            r.font.color.rgb = NAVY if (ci == 0 or total) else ORANGE; r.font.name = FONT
    rect(s, Inches(0.8), Inches(5.78), Inches(11.7), Inches(1.0), NAVY)
    text(s, Inches(1.05), Inches(5.86), Inches(11.3), Inches(0.85),
         [[("With NIC Karachi incubation:  ", 12, True, ORANGE),
           ("free office + zero utilities save ~PKR 80K to 120K/mo , that single saving dwarfs our entire cloud burn, redirected straight into provider acquisition.", 11, False, WHITE)],
          [("All 3 co-founders currently work deferred / equity-only. No salaries drawn pre-revenue.", 10, True, CL)]],
         space_after=3, line_spacing=1.04)
content_slide(10, "9. COST STRUCTURE  (CapEx vs OpEx)", "A lean, software-margin business", s10)

# =====================================================================
# 11 , FINANCIALS (with 3-year revenue table)
# =====================================================================
def s11(s):
    text(s, Inches(0.8), Inches(1.95), Inches(11.5), Inches(0.5),
         [[("Pre-revenue MVP today , high-margin economics as marketplace liquidity builds.", 15, True, ORANGE)]])
    cards = [("Revenue driver", "Completed jobs x commission per job. Grows with provider liquidity & repeat rate."),
             ("Gross margin", "Software margins , minimal cost per transaction once a city reaches density."),
             ("Burn", "Lean: cloud + APIs + acquisition. No inventory, no fleet, no salaries beyond core team.")]
    x = Inches(0.8)
    for head, body in cards:
        rect(s, x, Inches(2.7), Inches(3.8), Inches(1.6), WHITE, line=LITEOR)
        rect(s, x, Inches(2.7), Inches(3.8), Inches(0.12), ORANGE)
        text(s, x + Inches(0.25), Inches(2.95), Inches(3.3), Inches(1.3),
             [[(head, 14.5, True, NAVY)], [(body, 12, False, GRAY)]], space_after=6, line_spacing=1.12)
        x += Inches(4.0)
    # 3-year revenue table
    text(s, Inches(0.8), Inches(4.45), Inches(11.5), Inches(0.4),
         [[("Illustrative 3-year revenue trajectory:", 13, True, NAVY)]])
    rows, cols = 4, 5
    tbl = s.shapes.add_table(rows, cols, Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.6)).table
    tbl.columns[0].width = Inches(3.0)
    for c in range(1, cols):
        tbl.columns[c].width = Inches(2.175)
    hdr3 = ["", "Year 1 (Pilot)", "Year 2 (Karachi)", "Year 3 (Multi-city)", "Notes"]
    data3 = [
        ["Active providers", "40", "120", "175+", "Karachi mature + Lahore entry"],
        ["Monthly jobs", "480", "1,920", "4,375", "Providers x jobs/mo"],
        ["Annual revenue (PKR)", "~1.4M", "~8M", "~25M", "Commission-based, bottom-up"],
    ]
    for c in range(cols):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr3[c]; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    for ri, row in enumerate(data3, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.fill.solid()
            cell.fill.fore_color.rgb = CREAM if ri % 2 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val; r.font.size = Pt(10.5); r.font.bold = (ci == 0)
            r.font.color.rgb = NAVY if ci == 0 else ORANGE; r.font.name = FONT
    text(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
         [[("Unit economics: avg job ~PKR 600 \u2192 revenue/job = PKR 72 \u2192 break-even ~280 jobs/mo.  Pre-revenue , joining NIC to build professional financial models.", 10, True, LGRAY)]])
content_slide(11, "10. FINANCIALS", "Lean burn, scalable margins", s11)

# =====================================================================
# 12 , KEY METRICS & PROJECTIONS (timeline from pilot launch)
# =====================================================================
def s12(s):
    text(s, Inches(0.8), Inches(1.9), Inches(11.6), Inches(0.6),
         [[("Bottom-up plan (illustrative, from pilot launch) , built from active providers x jobs x Rs. 72 commission/job.", 13.5, True, ORANGE)]])
    rows, cols = 6, 5
    left, top = Inches(0.8), Inches(2.55)
    w, h = Inches(11.7), Inches(3.0)
    tbl = s.shapes.add_table(rows, cols, left, top, w, h).table
    headers = ["Metric", "Month 3", "Month 6", "Month 9", "Month 12"]
    metrics = [
        ["Active providers", "40", "75", "120", "175"],
        ["Jobs / month", "480", "1,050", "1,920", "3,150"],
        ["Commission / job", "Rs. 72", "Rs. 72", "Rs. 72", "Rs. 72"],
        ["Monthly revenue (PKR)", "~35K", "~76K", "~138K", "~227K"],
        ["Customer CAC (PKR)", "500", "350", "230", "150"],
    ]
    tbl.columns[0].width = Inches(3.1)
    for c in range(1, cols):
        tbl.columns[c].width = Inches(2.15)
    for c in range(cols):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = headers[c]; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for ri, mrow in enumerate(metrics, start=1):
        for ci, val in enumerate(mrow):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CREAM if ri % 2 else WHITE
            if ci == 0:
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.bold = (ci == 0)
            r.font.color.rgb = NAVY if ci == 0 else ORANGE
            r.font.name = FONT
    text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.1),
         [[("LTV ~Rs. 1,300 ", 11.5, True, ORANGE), ("(est. 6 jobs/yr x Rs. 72 x ~3-yr retention, based on Urban Company disclosed repeat rates).    ", 11.5, False, GRAY),
           ("Customer CAC: Rs. 500 \u2192 Rs. 150 ", 11.5, True, NAVY), ("as referral loops compound.", 11.5, False, GRAY)],
          [("Provider acquisition (field agent + signups) budgeted separately. Y1 revenue ~Rs. 1.4M (illustrative bottom-up) , validated in pilot, refined with NIC.", 10.5, True, LGRAY)]],
         space_after=4, line_spacing=1.1)
content_slide(12, "11. KEY METRICS & PROJECTIONS", "Bottom-up economics (from pilot launch)", s12)

# =====================================================================
# 13 , TECH STACK (simplified + visual roadmap)
# =====================================================================
def s13(s):
    cols = [
        ("FRONTEND", ["Next.js 15 + React + Tailwind", "PWA (offline ready), Capacitor (Android)", "Google Maps live tracking UI"]),
        ("BACKEND", ["Supabase Postgres + PostGIS", "RLS + SECURITY DEFINER RPCs", "Realtime (broadcast + postgres_changes)"]),
        ("SERVICES", ["Twilio Verify , phone OTP", "Firebase Cloud Messaging , push", "Google Maps , geocoding & maps"]),
    ]
    x = Inches(0.8)
    for title_, items in cols:
        rect(s, x, Inches(2.0), Inches(3.8), Inches(2.6), WHITE, line=LITEOR)
        rect(s, x, Inches(2.0), Inches(3.8), Inches(0.5), NAVY)
        text(s, x, Inches(2.05), Inches(3.8), Inches(0.45),
             [[(title_, 13, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bullets(s, x + Inches(0.2), Inches(2.65), Inches(3.5), Inches(1.9), items, size=11.5, gap=10)
        x += Inches(4.0)
    # Visual roadmap timeline
    text(s, Inches(0.8), Inches(4.85), Inches(11.6), Inches(0.4),
         [[("Roadmap", 14, True, NAVY)]])
    rect(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.08), ORANGE)
    phases = [
        ("NOW", "Pilot", "Karachi beta\n40 providers\nPayment webhooks"),
        ("Q2 to Q3", "City Scale", "175+ providers\nMobile apps\nKYC & CNIC checks"),
        ("Q4+", "Multi-City", "Lahore entry\nB2B2C enterprise\nProvider financing"),
    ]
    px = Inches(1.2)
    for label, title, desc in phases:
        oval(s, px + Inches(0.7), Inches(5.15), Inches(0.35), Inches(0.35), ORANGE)
        text(s, px, Inches(5.55), Inches(2.8), Inches(1.3),
             [[(label, 10, True, ORANGE)], [(title, 13, True, NAVY)], [(desc, 10, False, GRAY)]],
             space_after=2, line_spacing=1.05)
        px += Inches(3.8)
content_slide(13, "12. TECH STACK & ROADMAP", "Engineered to scale", s13)

# =====================================================================
# 14 , TEAM (single slide, real bios for all co-founders)
# =====================================================================
def s14(s):
    text(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.5),
         [[("A technical founder, backed by product, finance, growth & design.", 15, True, ORANGE)]])
    members = [
        ("MAR", "Marjan Ahmed", "Founder & CEO", "The Hacker / CTO",
         "Full-stack engineer , built all of USTAZ solo at 17. GIAIC Senior \u00b7 Agentic-AI engineer at Filion Capital \u00b7 AI-finance hackathon winner. 17-year-old who shipped a production-grade platform with 8+ core systems."),
        ("MAS", "Masood Alam", "Product Director & Financial Consultant", "The Hustler",
         "Leads product direction and the commission / wallet financial model. Background in operations and financial analysis. Owns unit economics, pricing strategy, and the prepaid wallet system design."),
        ("MSA", "Muhammad Sufyan Ahmed", "CMO & Designer", "Growth & Brand",
         "Owns brand identity, growth marketing strategy, and the product design system & UX. Drives customer acquisition channels and the referral loop architecture."),
    ]
    x = Inches(0.8)
    for initials, name, role, arche, desc in members:
        rect(s, x, Inches(2.35), Inches(3.8), Inches(3.7), NAVY)
        rect(s, x, Inches(2.35), Inches(3.8), Inches(0.12), ORANGE)
        oval(s, x + Inches(1.5), Inches(2.65), Inches(0.8), Inches(0.8), ORANGE)
        text(s, x + Inches(1.5), Inches(2.68), Inches(0.8), Inches(0.72),
             [[(initials, 15, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(0.25), Inches(3.6), Inches(3.35), Inches(2.4),
             [[(name, 14.5, True, WHITE)],
              [(role, 11, True, ORANGE)],
              [(arche + " , NIC archetype", 9.5, True, CL)],
              [(desc, 10.5, False, CL)]],
             space_after=4, line_spacing=1.06)
        x += Inches(4.0)
    text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.62),
         [[("Commitment:  ", 12, True, ORANGE),
           ("all 3 co-founders work full-time, equity-vested (4 year cliff). No salaries drawn pre-revenue. Staying lean: batch-verify through trade guilds (20 to 30 tradesmen at a time) to scale supply without drowning in manual checks.", 10.5, False, GRAY)]],
         line_spacing=1.04)
content_slide(14, "13. FOUNDERS & CORE TEAM", "Already built. Ready to scale.", s14)

# =====================================================================
# 15 , THE ASK (NEW , critical gap)
# =====================================================================
def s15(s):
    text(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.6),
         [[("What we need from NIC Karachi incubation:", 18, True, NAVY)]])
    asks = [
        ("Legal & SECP", "Registration support to formalize as a Pvt Ltd , required before any investment or enterprise partnerships can proceed."),
        ("Financial Modeling", "Mentorship to build professional P&L models, validate unit economics with real pilot data, and prepare for seed stage fundraising."),
        ("Karachi Enterprise Network", "Introductions to residential society management companies and B2B2C channels that unlock 200+ households per partnership."),
        ("Sharia Advisory", "Formal ujrah sign-off for the commission model , critical for trust in the Pakistani market and provider buy-in."),
        ("Provider Acquisition", "Access to NIC's network for trade guild partnerships and batch-onboarding channels across Karachi."),
        ("Market Validation", "Structured pilot framework to run 50+ real jobs, measure retention, and refine CAC before scaling."),
    ]
    y = Inches(2.6)
    for title, desc in asks:
        rect(s, Inches(0.8), y, Inches(0.12), Inches(0.7), ORANGE)
        text(s, Inches(1.1), y, Inches(2.5), Inches(0.7),
             [[(title, 13, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(3.6), y, Inches(9.0), Inches(0.7),
             [[(desc, 12, False, GRAY)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += Inches(0.78)
    text(s, Inches(0.8), Inches(6.3), Inches(11.6), Inches(0.5),
         [[("We are not just building an app , we are building Pakistan's trusted home-services infrastructure. NIC's ecosystem is the catalyst to get from MVP to market.", 12, True, ORANGE)]])
content_slide(15, "14. THE ASK", "What we need from NIC Karachi", s15)

# =====================================================================
# 16 , LONG-TERM VISION & EXIT
# =====================================================================
def s16(s):
    text(s, Inches(0.8), Inches(1.95), Inches(11.5), Inches(0.5),
         [[("Exit path", 16, True, ORANGE)]])
    bullets(s, Inches(0.8), Inches(2.5), Inches(11.6), Inches(2), [
        ("Primary , strategic acquisition: ", "by a super-app, classifieds leader, telco, or bank scaling into home services."),
        ("Likely acquirers: ", "Jazz / PTCL (telco super-app), OLX Pakistan (classifieds to services), Careem Pakistan, HBL / Meezan (bank-led services)."),
        ("Secondary , regional roll-up: ", "dominate Pakistan, then expand to comparable South Asian markets."),
    ], size=13.5, gap=10)
    text(s, Inches(0.8), Inches(4.35), Inches(11.5), Inches(0.5),
         [[("Corporate structure & governance", 16, True, ORANGE)]])
    bullets(s, Inches(0.8), Inches(4.9), Inches(11.6), Inches(2), [
        ("SECP registration: ", "filing as Pvt Ltd within 3 to 6 months of incubation, with NIC legal mentorship. May register as sole-proprietorship immediately to show intent."),
        ("Equity: ", "3 co-founders on an agreed split (~60/20/20); 4 year vesting + anti dilution & exit clauses formalized with NIC legal support pre-investment."),
    ], size=13.5, gap=10)
content_slide(16, "15. LONG-TERM VISION & EXIT STRATEGY", "Building Pakistan's default home-services brand", s16)

# =====================================================================
# 17 , CONTACT / THANK YOU
# =====================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, WHITE)
rect(s, 0, 0, Inches(4.6), EMU_H, NAVY)
oval(s, Inches(-1.0), Inches(5.4), Inches(2.6), Inches(2.6), ORANGE)
text(s, Inches(0.6), Inches(2.7), Inches(3.6), Inches(1.4),
     [[("USTAZ", 44, True, WHITE)], [("Trusted pros, on demand.", 14, True, ORANGE)]], space_after=8)
text(s, Inches(5.2), Inches(1.4), Inches(7.5), Inches(0.8), [[("Thank you", 40, True, NAVY)]])
text(s, Inches(5.2), Inches(2.5), Inches(7.5), Inches(0.5), [[("CONTACT", 13, True, ORANGE)]])
contact = [
    ("Website", "ustaz-bice.vercel.app"),
    ("Email", "marjanahmed.dev@gmail.com"),
    ("Founder", "Marjan Ahmed , Founder & CEO"),
    ("Location", "Karachi, Pakistan"),
]
y = Inches(3.15)
for label, val in contact:
    rect(s, Inches(5.2), y + Inches(0.07), Inches(0.12), Inches(0.4), ORANGE)
    text(s, Inches(5.45), y, Inches(2.0), Inches(0.5), [[(label, 13, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.4), y, Inches(5.2), Inches(0.5), [[(val, 13, False, GRAY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.62)
text(s, Inches(5.2), Inches(6.2), Inches(7.5), Inches(0.5),
     [[("Applying for NIC Karachi incubation", 11, False, LGRAY)]])

prs.save("USTAZ-Pitch-Deck.pptx")
print("Saved USTAZ-Pitch-Deck.pptx with", len(prs.slides._sldIdLst), "slides")
